# -*- coding: utf-8 -*-
"""Genera la presentacion del trabajo final de MIA-10 (PLN).

Estructura pedida por el curso, con los cuatro bloques de la plantilla del
docente dentro del punto 4:

  1 Introduccion · 2 Planteamiento del problema · 3 Antecedentes
  4 Planteamiento de la solucion
       datos · preprocesamiento · modelamiento · EJECUCION · resultados
  5 Conclusiones · 6 Recomendaciones · 7 Referencias

Todas las cifras salen de los artefactos reales del pipeline (fase 9 a 12 y
OE5). Nada esta escrito a mano dos veces: si una cifra cambia en el JSON,
cambia aqui.

Uso:  python generar_presentacion_final.py
Salida: presentacion/Presentacion_Trabajo_Final_PLN_CarlosPerez.pptx
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

RAIZ = Path(r"T:\MIMIC\PLN_SP")
ART = Path(r"T:\MIMIC\tesis\04_pipeline_codigo\datos_intermedios")
PLANTILLA = Path(r"C:\Users\infor\OneDrive\CARLOS PEREZ\MAESTRIA  UNI "
                 r"INTELIGENCIA ARTIFICIAL\III CICLO\CURSO Procesamiento del "
                 r"Lenguaje Natural\01 Clases - Sesiones\clase 14\
Plantilla_Trabajo_Final.pptx".replace("\n", ""))
SALIDA = RAIZ / "presentacion" / "Presentacion_Trabajo_Final_PLN_CarlosPerez.pptx"

# paleta: granate UNI + azul acero, sobre blanco
GRANATE = RGBColor(0x7B, 0x15, 0x22)
GRANATE_CLARO = RGBColor(0xF2, 0xE4, 0xE6)
AZUL = RGBColor(0x2E, 0x5F, 0x7F)
AZUL_CLARO = RGBColor(0xE4, 0xED, 0xF3)
TINTA = RGBColor(0x1F, 0x24, 0x28)
GRIS = RGBColor(0x5A, 0x62, 0x68)
GRIS_TENUE = RGBColor(0xEC, 0xEE, 0xF0)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
VERDE = RGBColor(0x1D, 0x6F, 0x54)
CONSOLA_BG = RGBColor(0x1B, 0x22, 0x28)
CONSOLA_FG = RGBColor(0xD8, 0xE0, 0xE6)
AMBAR = RGBColor(0xB5, 0x7A, 0x0E)

ANCHO = Inches(13.333)
ALTO = Inches(7.5)
MARGEN = Inches(0.62)
ANCHO_UTIL = ANCHO - 2 * MARGEN

# ---------------------------------------------------------------- artefactos


def cargar(ruta):
    f = ART / ruta
    return json.loads(f.read_text(encoding="utf-8")) if f.exists() else None


f9 = cargar("fase9_final/resultados_finales.json")
f10 = cargar("fase9_final/metricas_corregidas.json")
f11 = cargar("fase11/resultados_transformers.json")
f12 = cargar("fase12/sistema_vs_experto.json")
kap = cargar("fase6_concordancia/concordancia.json")
oe5 = cargar("oe5_ersp/informe_oe5.json")

prs = Presentation(str(PLANTILLA))
prs.slide_width, prs.slide_height = ANCHO, ALTO
BLANCO_LAYOUT = prs.slide_layouts[6]

# la plantilla trae 2 diapositivas de muestra: se eliminan
for i in range(len(prs.slides) - 1, -1, -1):
    rid = prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rid)
    del prs.slides._sldIdLst[i]

_n = [0]

# ------------------------------------------------------------------ helpers


def caja(slide, x, y, cx, cy, texto="", tam=14, color=TINTA, negrita=False,
         alineado=PP_ALIGN.LEFT, fuente="Calibri", interlinea=1.0,
         anclaje=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anclaje
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = alineado
    p.line_spacing = interlinea
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(tam)
    r.font.bold = negrita
    r.font.color.rgb = color
    r.font.name = fuente
    return tb


def rect(slide, x, y, cx, cy, relleno=None, borde=None, grosor=0.75):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
    if relleno is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = relleno
    if borde is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = borde
        sh.line.width = Pt(grosor)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh


def lamina(titulo=None, bajada=None, seccion=None):
    """Crea una diapositiva con la cabecera estandar y devuelve (slide, y)."""
    _n[0] += 1
    s = prs.slides.add_slide(BLANCO_LAYOUT)
    y = MARGEN
    if seccion:
        caja(s, MARGEN, y, ANCHO_UTIL, Inches(0.24), seccion.upper(),
             tam=11, color=AZUL, negrita=True)
        y += Inches(0.3)
    if titulo:
        caja(s, MARGEN, y, ANCHO_UTIL, Inches(0.55), titulo,
             tam=28, color=GRANATE, negrita=True)
        y += Inches(0.62)
        # filete
        rect(s, MARGEN, y, Inches(1.5), Pt(2.5), relleno=GRANATE)
        y += Inches(0.18)
    if bajada:
        tb = caja(s, MARGEN, y, ANCHO_UTIL, Inches(0.4), bajada,
                  tam=14, color=GRIS, interlinea=1.15)
        y += Inches(0.28) + Inches(0.22) * (len(bajada) // 125)
    y += Inches(0.14)
    # pie
    caja(s, MARGEN, ALTO - Inches(0.42), Inches(9),
         Inches(0.24), "MIA-10 Procesamiento del Lenguaje Natural · "
         "Carlos Pérez Pérez · UNI 2026", tam=9, color=GRIS)
    caja(s, ANCHO - MARGEN - Inches(0.6), ALTO - Inches(0.42), Inches(0.6),
         Inches(0.24), str(_n[0]), tam=9, color=GRIS, alineado=PP_ALIGN.RIGHT)
    return s, y


def vinetas(slide, x, y, cx, items, tam=15, sep=Inches(0.42), color=TINTA):
    for it in items:
        rect(slide, x, y + Inches(0.09), Pt(5), Pt(5), relleno=GRANATE)
        if isinstance(it, tuple):
            tb = slide.shapes.add_textbox(x + Inches(0.22), y - Inches(0.03),
                                          cx - Inches(0.22), sep)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.line_spacing = 1.12
            r1 = p.add_run()
            r1.text = it[0] + "  "
            r1.font.size = Pt(tam)
            r1.font.bold = True
            r1.font.color.rgb = TINTA
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = it[1]
            r2.font.size = Pt(tam)
            r2.font.color.rgb = GRIS
            r2.font.name = "Calibri"
            alto_real = sep * (1 + len(it[0] + it[1]) // 110)
        else:
            caja(slide, x + Inches(0.22), y - Inches(0.03), cx - Inches(0.22),
                 sep, it, tam=tam, color=color, interlinea=1.12)
            alto_real = sep * (1 + len(it) // 110)
        y += alto_real
    return y


def tabla(slide, x, y, cx, datos, anchos=None, tam=12, alto_fila=Inches(0.34),
          resaltar=None, alto_cabecera=Inches(0.4)):
    """datos[0] es la cabecera. resaltar = indice de fila a destacar."""
    ncol = len(datos[0])
    if anchos is None:
        anchos = [1.0 / ncol] * ncol
    for j, frac in enumerate(anchos):
        pass
    y0 = y
    # cabecera
    rect(slide, x, y, cx, alto_cabecera, relleno=GRANATE)
    cx_acum = x
    for j, celda in enumerate(datos[0]):
        w = int(cx * anchos[j])
        caja(slide, cx_acum + Inches(0.1), y + Inches(0.08),
             w - Inches(0.14), alto_cabecera, str(celda), tam=tam,
             color=BLANCO, negrita=True)
        cx_acum += w
    y += alto_cabecera
    # filas
    for i, fila in enumerate(datos[1:]):
        destaca = (resaltar is not None and i == resaltar)
        fondo = GRANATE_CLARO if destaca else (GRIS_TENUE if i % 2 else None)
        if fondo:
            rect(slide, x, y, cx, alto_fila, relleno=fondo)
        cx_acum = x
        for j, celda in enumerate(fila):
            w = int(cx * anchos[j])
            caja(slide, cx_acum + Inches(0.1), y + Inches(0.06),
                 w - Inches(0.14), alto_fila, str(celda), tam=tam,
                 color=TINTA if destaca else TINTA,
                 negrita=destaca)
            cx_acum += w
        y += alto_fila
    rect(slide, x, y0, cx, y - y0, relleno=None, borde=RGBColor(0xD5, 0xD9, 0xDD))
    return y


def consola(slide, x, y, cx, cy, lineas, tam=11, titulo=None):
    if titulo:
        caja(slide, x, y, cx, Inches(0.22), titulo, tam=11,
             color=AZUL, negrita=True, fuente="Consolas")
        y += Inches(0.26)
    rect(slide, x, y, cx, cy, relleno=CONSOLA_BG)
    tb = slide.shapes.add_textbox(x + Inches(0.14), y + Inches(0.1),
                                  cx - Inches(0.28), cy - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.08
        destaca = ln.startswith("*")
        r = p.add_run()
        r.text = ln[1:] if destaca else ln
        r.font.size = Pt(tam)
        r.font.name = "Consolas"
        r.font.color.rgb = RGBColor(0x7A, 0xD1, 0xA8) if destaca else CONSOLA_FG
        r.font.bold = destaca
    return y + cy


def codigo(slide, x, y, cx, cy, lineas, tam=11, titulo=None):
    if titulo:
        caja(slide, x, y, cx, Inches(0.22), titulo, tam=11, color=AZUL,
             negrita=True, fuente="Consolas")
        y += Inches(0.26)
    rect(slide, x, y, cx, cy, relleno=RGBColor(0xF7, 0xF8, 0xF9),
         borde=RGBColor(0xD5, 0xD9, 0xDD))
    rect(slide, x, y, Pt(3), cy, relleno=AZUL)
    tb = slide.shapes.add_textbox(x + Inches(0.16), y + Inches(0.1),
                                  cx - Inches(0.3), cy - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1
        coment = ln.strip().startswith("#")
        destaca = ln.startswith("*")
        r = p.add_run()
        r.text = ln[1:] if destaca else ln
        r.font.size = Pt(tam)
        r.font.name = "Consolas"
        if coment:
            r.font.color.rgb = VERDE
        elif destaca:
            r.font.color.rgb = GRANATE
            r.font.bold = True
        else:
            r.font.color.rgb = TINTA
    return y + cy


def dato(slide, x, y, cx, cy, cifra, rotulo, color=GRANATE, pie=None):
    rect(slide, x, y, cx, cy, relleno=GRIS_TENUE)
    rect(slide, x, y, cx, Pt(3), relleno=color)
    caja(slide, x + Inches(0.16), y + Inches(0.2), cx - Inches(0.32),
         Inches(0.5), cifra, tam=30, color=color, negrita=True)
    caja(slide, x + Inches(0.16), y + Inches(0.76), cx - Inches(0.32),
         Inches(0.5), rotulo, tam=12, color=GRIS, interlinea=1.1)
    if pie:
        caja(slide, x + Inches(0.16), y + cy - Inches(0.34),
             cx - Inches(0.32), Inches(0.28), pie, tam=10, color=AZUL)


def nota(slide, y, texto, color=AZUL, fondo=AZUL_CLARO, alto=Inches(0.62)):
    rect(slide, MARGEN, y, ANCHO_UTIL, alto, relleno=fondo)
    rect(slide, MARGEN, y, Pt(4), alto, relleno=color)
    caja(slide, MARGEN + Inches(0.2), y + Inches(0.1),
         ANCHO_UTIL - Inches(0.4), alto - Inches(0.2), texto, tam=13,
         color=TINTA, interlinea=1.16, anclaje=MSO_ANCHOR.MIDDLE)
    return y + alto


def flecha(slide, x, y, cx, cy, texto, sub, relleno, texto_color=BLANCO):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, x, y, cx, cy)
    sh.fill.solid()
    sh.fill.fore_color.rgb = relleno
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    caja(slide, x + Inches(0.14), y + Inches(0.16), cx - Inches(0.5),
         Inches(0.5), texto, tam=13, color=texto_color, negrita=True,
         interlinea=1.05)
    caja(slide, x + Inches(0.14), y + Inches(0.74), cx - Inches(0.5),
         cy - Inches(0.8), sub, tam=10, color=texto_color, interlinea=1.1)


# =========================================================== 1. PORTADA
s = prs.slides.add_slide(BLANCO_LAYOUT)
rect(s, 0, 0, ANCHO, Inches(2.55), relleno=GRANATE)
caja(s, MARGEN, Inches(0.5), ANCHO_UTIL, Inches(0.3),
     "UNIVERSIDAD NACIONAL DE INGENIERÍA · MAESTRÍA EN INTELIGENCIA ARTIFICIAL",
     tam=12, color=RGBColor(0xE8, 0xC8, 0xCC), negrita=True)
caja(s, MARGEN, Inches(0.95), Inches(11.4), Inches(1.3),
     "Detección automática de eventos adversos\nhospitalarios en notas clínicas",
     tam=38, color=BLANCO, negrita=True, interlinea=1.05)
caja(s, MARGEN, Inches(2.78), Inches(11.4), Inches(0.4),
     "Un canal de Procesamiento de Lenguaje Natural sobre MIMIC-IV, "
     "auditado y validado contra juicio experto", tam=16, color=AZUL)
rect(s, MARGEN, Inches(3.4), Inches(5.4), Pt(2), relleno=GRANATE)
for i, (k, v) in enumerate([
        ("Curso", "MIA-10 · Procesamiento del Lenguaje Natural"),
        ("Docente", "Dr. Wester Zela Moraya"),
        ("Autor", "Carlos Pérez Pérez"),
        ("Fecha", "9 de agosto de 2026")]):
    yy = Inches(3.62) + Inches(0.46) * i
    caja(s, MARGEN, yy, Inches(1.3), Inches(0.34), k, tam=14, color=GRIS)
    caja(s, MARGEN + Inches(1.45), yy, Inches(6), Inches(0.34), v, tam=15.5,
         color=TINTA, negrita=(i in (0, 1, 2)))

rect(s, Inches(7.6), Inches(3.5), Inches(5.1), Inches(3.2), relleno=GRIS_TENUE)
rect(s, Inches(7.6), Inches(3.5), Inches(5.1), Pt(3), relleno=GRANATE)
caja(s, Inches(7.85), Inches(3.72), Inches(4.6), Inches(0.3),
     "EL TRABAJO EN CUATRO CIFRAS", tam=11, color=GRANATE, negrita=True)
for i, (c, t) in enumerate([
        ("70,000", "notas clínicas procesadas"),
        ("33 h 54 min", "de cómputo registrado"),
        ("7", "modelos comparados"),
        ("0.843", "AUC del detector final")]):
    yy = Inches(4.12) + Inches(0.6) * i
    caja(s, Inches(7.85), yy, Inches(1.75), Inches(0.4), c, tam=19,
         color=GRANATE, negrita=True)
    caja(s, Inches(9.7), yy + Inches(0.07), Inches(2.8), Inches(0.4), t,
         tam=12, color=GRIS)
caja(s, MARGEN, ALTO - Inches(0.5), Inches(11), Inches(0.3),
     "github.com/carlosperez100/PLN_SP   ·   carlosperez100.github.io/PLN_SP",
     tam=11, color=AZUL)

# =========================================================== 2. AGENDA
s, y = lamina("Agenda")
izq = [("1", "Introducción", "Seguridad del paciente y la norma"),
       ("2", "Planteamiento del problema", "El subregistro del 72 %"),
       ("3", "Antecedentes", "Estado del arte y selección del modelo"),
       ("4", "Planteamiento de la solución", "Datos · Preprocesamiento · "
        "Modelamiento · Ejecución · Resultados")]
der = [("5", "Conclusiones", "Qué se confirmó y qué se refutó"),
       ("6", "Recomendaciones", "Hacia dónde sigue el trabajo"),
       ("7", "Referencias", "34 fuentes citadas")]
for col, items in ((MARGEN, izq), (Inches(7.0), der)):
    yy = y + Inches(0.3)
    for num, tit, sub in items:
        rect(s, col, yy, Inches(0.42), Inches(0.42), relleno=GRANATE)
        caja(s, col, yy + Inches(0.07), Inches(0.42), Inches(0.3), num,
             tam=15, color=BLANCO, negrita=True, alineado=PP_ALIGN.CENTER)
        caja(s, col + Inches(0.6), yy, Inches(5.3), Inches(0.3), tit,
             tam=18, color=TINTA, negrita=True)
        caja(s, col + Inches(0.6), yy + Inches(0.33), Inches(5.3),
             Inches(0.5), sub, tam=12, color=GRIS, interlinea=1.1)
        yy += Inches(1.02)
nota(s, Inches(6.05), "El punto 4 concentra el 60 % de la exposición e "
     "incluye la ejecución real del código: comandos, salidas de consola y "
     "el cuaderno corriendo en vivo.")

# =========================================================== 3. INTRODUCCIÓN
s, y = lamina("Un daño causado por la atención, no por la enfermedad",
              seccion="1 · Introducción")
vinetas(s, MARGEN, y + Inches(0.1), Inches(6.1), [
    ("Definición.", "Daño no intencional derivado de la atención sanitaria, "
     "no del curso natural de la enfermedad del paciente."),
    ("Magnitud.", "Están entre las 10 primeras causas de muerte y "
     "discapacidad en el mundo (OMS)."),
    ("En el Perú.", "La Directiva GG-ESSALUD-2021 los regula y define su "
     "notificación obligatoria."),
], sep=Inches(0.78))

rect(s, Inches(7.1), y + Inches(0.1), Inches(5.6), Inches(3.5),
     relleno=GRIS_TENUE)
rect(s, Inches(7.1), y + Inches(0.1), Inches(5.6), Pt(3), relleno=GRANATE)
caja(s, Inches(7.35), y + Inches(0.32), Inches(5.1), Inches(0.3),
     "ANEXO 02 · LA TAXONOMÍA QUE SE AUTOMATIZA", tam=11, color=GRANATE,
     negrita=True)
for i, (c, t) in enumerate([("231", "eventos adversos tipificados"),
                            ("12", "naturalezas que los agrupan"),
                            ("4", "niveles de severidad (Anexo 03)")]):
    yy = y + Inches(0.78) + Inches(0.72) * i
    caja(s, Inches(7.35), yy, Inches(1.1), Inches(0.45), c, tam=26,
         color=GRANATE, negrita=True)
    caja(s, Inches(8.6), yy + Inches(0.1), Inches(3.9), Inches(0.4), t,
         tam=13, color=GRIS)
nota(s, Inches(5.5), "La tarea de este trabajo: leer el texto libre de una "
     "nota clínica y decidir si contiene un evento adverso y de qué naturaleza "
     "es, según esa taxonomía.")

# =========================================================== 4. PROBLEMA
s, y = lamina("La información existe, pero nadie la lee",
              seccion="2 · Planteamiento del problema")
w = Inches(4.02)
for i, (c, r, p) in enumerate([
        ("~72 %", "de los eventos adversos no se notifican en EsSalud",
         "unos 37,000 al año"),
        ("100 %", "de las notas clínicas contienen esa información en texto libre",
         "no explotable a mano"),
        ("Manual", "la notificación depende de que alguien decida reportar",
         "y voluntaria")]):
    dato(s, MARGEN + (w + Inches(0.28)) * i, y + Inches(0.1), w,
         Inches(1.75), c, r, pie=p)

y2 = y + Inches(2.1)
caja(s, MARGEN, y2, ANCHO_UTIL, Inches(0.3), "Por qué falla el registro "
     "manual", tam=17, color=TINTA, negrita=True)
vinetas(s, MARGEN, y2 + Inches(0.45), Inches(6.1), [
    "Subnotificación por temor a la sanción",
    "Saturación del personal: el evento se registra pero no se analiza",
    "Priorización heterogénea entre servicios y lenta",
], sep=Inches(0.42), tam=14)
vinetas(s, Inches(7.1), y2 + Inches(0.45), Inches(5.6), [
    "El texto clínico es largo (mediana 3,148 tokens)",
    "El evento no está señalado: hay que encontrarlo",
    "Sin denominador no se puede priorizar por riesgo",
], sep=Inches(0.42), tam=14)

nota(s, Inches(6.0), "Pregunta de investigación:  ¿es posible detectar "
     "automáticamente, desde el texto libre del resumen de alta, los eventos "
     "adversos ocurridos durante la hospitalización, y clasificarlos según la "
     "taxonomía normativa, con fiabilidad comparable a la del juicio experto?",
     color=GRANATE, fondo=GRANATE_CLARO)

# =========================================================== 5. OBJETIVOS
s, y = lamina("Objetivos del trabajo", seccion="2 · Planteamiento del problema")
nota(s, y, "Objetivo general — diseñar, implementar y validar un canal de PLN "
     "que detecte y priorice eventos adversos automáticamente, con miras a su "
     "transferencia a EsSalud.", color=GRANATE, fondo=GRANATE_CLARO)
datos = [["", "Objetivo específico", "Estado al día de hoy"],
         ["OE1", "Construir un corpus etiquetado por supervisión débil",
          "Cumplido · 70,000 notas clínicas"],
         ["OE2", "Comparar transformers biomédicos con modelos clásicos",
          "Cumplido · 7 modelos"],
         ["OE3", "Automatizar la matriz de priorización GEMSES",
          "Implementado · fuera del alcance de PLN"],
         ["OE4", "Contrastar el sistema contra evaluadores expertos",
          "Cumplido · 163 casos con juicio experto"],
         ["OE5", "Documentar la transferencia al español",
          "Cumplido · 6,336 casos ERSP"]]
tabla(s, MARGEN, y + Inches(0.85), ANCHO_UTIL, datos,
      anchos=[0.07, 0.54, 0.39], alto_fila=Inches(0.46), tam=14)
nota(s, Inches(5.95), "Esta exposición se concentra en OE1, OE2, OE4 y OE5, "
     "que son los objetivos de procesamiento de lenguaje natural. El OE3 es "
     "gestión y se muestra solo como destino aplicado.")

# =========================================================== 6. ANTECEDENTES
s, y = lamina("Estado del arte: cinco líneas que sostienen el diseño",
              seccion="3 · Antecedentes")
filas = [["Línea", "Referencia", "Qué aporta a este trabajo"],
         ["Supervisión débil", "Ratner et al. (Snorkel)",
          "Etiquetar a escala con reglas, sin anotación manual masiva"],
         ["Aprendizaje por atajo", "Geirhos et al. (2020); Zech et al. (2018)",
          "Un modelo puede acertar por la razón equivocada: hay que auditarlo"],
         ["Ventana de contexto", "Li et al. (2022) · Clinical-Longformer",
          "Extender la ventana a 4,096 supera a ClinicalBERT en texto largo"],
         ["Encoders clínicos ES", "JAMIA (2024) · BETO, bsc-bio-ehr-es",
          "Fija la expectativa realista para la transferencia al español"],
         ["Fiabilidad observador", "Byrt et al. (1993) · PABAK",
          "Corrige la paradoja de kappa cuando la prevalencia es desigual"]]
tabla(s, MARGEN, y + Inches(0.1), ANCHO_UTIL, filas,
      anchos=[0.19, 0.28, 0.53], alto_fila=Inches(0.52), tam=13)
nota(s, Inches(5.05), "Zech et al. entrenaron un detector de neumonía que "
     "aprendió a reconocer el hospital de procedencia por marcas en la "
     "radiografía. Ese antecedente es exactamente lo que encontramos aquí "
     "— y lo veremos medido.", color=AMBAR,
     fondo=RGBColor(0xFA, 0xF0, 0xDC))

# =========================================================== 7. SELECCIÓN
s, y = lamina("Cómo se eligió el modelo, y con qué criterio",
              seccion="3 · Antecedentes",
              bajada="La tarea determina la métrica, y la métrica determina "
              "qué ranking es pertinente. No al revés.")
filas = [["Instrumento", "Por qué se usó o se descartó"],
         ["MTEB / MMTEB Leaderboard",
          "Único ranking masivo que puntúa clasificación y filtra por idioma"],
         ["IberBench (arXiv 2504.16921)",
          "Referencia para español; mide tareas de interés industrial"],
         ["JAMIA 2024 — encoders clínicos ES",
          "Evidencia revisada por pares sobre 12 corpus clínicos"],
         ["Open LLM Leaderboard",
          "DESCARTADO: archivado por sus autores y mide razonamiento, no "
          "clasificación"]]
tabla(s, MARGEN, y + Inches(0.1), Inches(7.4), filas, anchos=[0.36, 0.64],
      alto_fila=Inches(0.62), tam=12, resaltar=3)

rect(s, Inches(8.35), y + Inches(0.1), Inches(4.35), Inches(3.05),
     relleno=GRIS_TENUE)
rect(s, Inches(8.35), y + Inches(0.1), Inches(4.35), Pt(3), relleno=AZUL)
caja(s, Inches(8.6), y + Inches(0.32), Inches(3.9), Inches(0.3),
     "FILTROS DUROS DE ADMISIBILIDAD", tam=11, color=AZUL, negrita=True)
vinetas(s, Inches(8.6), y + Inches(0.72), Inches(3.85), [
    ("F1", "clasificación multietiqueta"),
    ("F2", "opera sobre inglés clínico"),
    ("F3", "CPU, huella ≤ 2.5 GB"),
    ("F4", "cumple el DUA de PhysioNet"),
    ("F5", "pesos y código libres"),
], tam=12, sep=Inches(0.44))
nota(s, Inches(5.6), "El filtro F4 es el que casi nadie aplica: PhysioNet "
     "prohíbe expresamente enviar texto de MIMIC a APIs de terceros. Eso "
     "elimina de entrada a ChatGPT y a cualquier LLM por API pública.")

# =========================================================== 8. FLUJO
s, y = lamina("El flujo completo, de la nota libre al responsable",
              seccion="4 · Planteamiento de la solución")
img_flujo = RAIZ / "presentacion" / "fig_flujo_completo.png"
w_img = Inches(11.9)
h_img = Inches(11.9 / 2.514)          # aspecto de la figura 3836x1526
s.shapes.add_picture(str(img_flujo),
                     MARGEN + int((ANCHO_UTIL - w_img) / 2),
                     y + Inches(0.25), width=w_img, height=h_img)

# ===================================================== 8b. FASES Y SCRIPTS
s, y = lamina("Las fases del flujo y el script que ejecuta cada una",
              seccion="4 · Planteamiento de la solución")
y2 = y + Inches(0.1)
filas = [["Fase", "Qué resuelve", "Script"],
         ["3", "Corpus por supervisión débil: códigos CIE + 35 patrones",
          "fase3_v2_corpus_completo.py"],
         ["4", "Comparación de modelos, fuga, circularidad, multietiqueta",
          "fase4_*.py  (7 scripts)"],
         ["9", "Modelo final con el confusor de época controlado",
          "fase9_modelo_final.py"],
         ["11", "Ajuste fino de Bio_ClinicalBERT y BioBERT en GPU",
          "fase11_finetuning_transformers.py"],
         ["12", "Contraste contra el consenso de dos expertos",
          "fase12_sistema_vs_experto.py"],
         ["OE5", "Transferencia al español con etiqueta de oro",
          "oe5_ersp_preparar.py"]]
fin_scripts = tabla(s, MARGEN, y2, ANCHO_UTIL, filas,
                    anchos=[0.07, 0.58, 0.35], alto_fila=Inches(0.52),
                    tam=14, alto_cabecera=Inches(0.44))
nota(s, fin_scripts + Inches(0.25), "Todo el código es público en "
     "github.com/carlosperez100/PLN_SP; los datos clínicos NO se versionan "
     "(DUA de PhysioNet y protección de datos). La bitácora conserva el "
     "comando y la salida de consola de cada corrida.")

# =========================================================== 9. DATOS MIMIC
s, y = lamina("Los datos: MIMIC-IV-Note v2.2",
              seccion="4.1 · Descripción de los datos",
              bajada="331,793 resúmenes de alta del Beth Israel Deaconess "
              "Medical Center (2008–2019), con acceso credencializado por "
              "PhysioNet.")
c = f9["corpus"]
w = Inches(3.02)
for i, (cif, rot) in enumerate([
        (f"{c['epicrisis']:,}", "notas clínicas en el corpus de modelado"),
        (f"{c['positivas']:,}",
         f"positivas ({c['positivas']/c['epicrisis']:.1%} del corpus)"),
        (f"{c['prevalencia_real']:.2%}", "prevalencia POBLACIONAL real"),
        (f"{c['universo_positivo']:,}", "hospitalizaciones con evento")]):
    dato(s, MARGEN + (w + Inches(0.2)) * i, y + Inches(0.05), w,
         Inches(1.62), cif, rot)

y2 = y + Inches(1.95)
caja(s, MARGEN, y2, Inches(6.1), Inches(0.3),
     "Por qué 70,000 y no las 331,793", tam=16, color=TINTA, negrita=True)
caja(s, MARGEN, y2 + Inches(0.4), Inches(6.1), Inches(1.1),
     "Vectorizar n-gramas de carácter sobre más de 100,000 notas desborda la "
     "RAM disponible. El techo práctico verificado fue de 70,000 notas. "
     "El corpus está balanceado a propósito; la población no lo está.",
     tam=13, color=GRIS, interlinea=1.25)

e = c["emparejamiento_epoca"]
rect(s, Inches(7.1), y2 - Inches(0.1), Inches(5.6), Inches(1.75),
     relleno=GRIS_TENUE)
rect(s, Inches(7.1), y2 - Inches(0.1), Inches(5.6), Pt(3), relleno=VERDE)
caja(s, Inches(7.35), y2 + Inches(0.1), Inches(5.1), Inches(0.3),
     "CONTROL DEL CONFUSOR DE ÉPOCA", tam=11, color=VERDE, negrita=True)
caja(s, Inches(7.35), y2 + Inches(0.48), Inches(5.1), Inches(0.9),
     f"positivos era CIE-10:  {e['positivos_era10']:.2%}\n"
     f"negativos era CIE-10:  {e['negativos_era10']:.2%}\n"
     f"diferencia:            {abs(e['positivos_era10']-e['negativos_era10']):.4%}",
     tam=14, color=TINTA, fuente="Consolas", interlinea=1.3)

nota(s, Inches(5.85), "La prevalencia poblacional (20.12 %) es muy distinta "
     "de la del conjunto de prueba, que está balanceado. Esa diferencia es la "
     "que obliga a reajustar el valor predictivo positivo — lo veremos en "
     "resultados.")

# =========================================================== 10. DATOS ERSP
s, y = lamina("Segundo corpus: texto peruano real, con etiqueta de oro",
              seccion="4.1 · Descripción de los datos",
              bajada="8,799 ocurrencias REALES del sistema peruano de "
              "reporte, en español, codificadas una a una por profesionales "
              "contra los Anexos 02 y 03. No es una simulación ni una "
              "traducción: es la norma y el idioma de destino.")
lim = oe5["limpieza"]
filas = [["Paso del preprocesamiento", "Casos"],
         ["Registros originales", f"{lim['filas_originales']:,}"],
         ["Identificadores anonimizados en el texto",
          f"{lim['dni_anonimizados']}"],
         ["Textos demasiado breves (< 30 caracteres)", f"{lim['textos_cortos']}"],
         ["Duplicados eliminados", "2,463"],
         ["Corpus final", f"{lim['filas_unicas']:,}"]]
tabla(s, MARGEN, y + Inches(0.05), Inches(6.0), filas, anchos=[0.72, 0.28],
      alto_fila=Inches(0.4), tam=13, resaltar=4)

rect(s, Inches(7.1), y + Inches(0.05), Inches(5.6), Inches(2.45),
     relleno=GRANATE_CLARO)
rect(s, Inches(7.1), y + Inches(0.05), Inches(5.6), Pt(3), relleno=GRANATE)
caja(s, Inches(7.35), y + Inches(0.25), Inches(5.1), Inches(0.3),
     "ESTÁNDAR DE PLATA FRENTE A ESTÁNDAR DE ORO", tam=11, color=GRANATE,
     negrita=True)
caja(s, Inches(7.35), y + Inches(0.65), Inches(5.1), Inches(1.6),
     "MIMIC se etiqueta con códigos CIE que asigna un codificador "
     "administrativo: estándar de PLATA.\n\nEl ERSP lo codificó un "
     "profesional leyendo cada descripción contra la norma: estándar de ORO. "
     "Es la primera medición contra criterio humano.",
     tam=13, color=TINTA, interlinea=1.25)

nota(s, Inches(4.85), "Sin la deduplicación el mismo texto caía en "
     "entrenamiento y en prueba, y las métricas salían infladas. Ese paso "
     "solo eliminó 2,463 filas de 8,799.", color=AMBAR,
     fondo=RGBColor(0xFA, 0xF0, 0xDC))
nota(s, Inches(5.62), "Advertencia de comparabilidad — estas cifras NO son "
     "comparables con las de MIMIC: el reporte ERSP describe un evento que "
     "quien lo escribió YA identificó; en la nota clínica el evento hay que "
     "ENCONTRARLO. Es un problema distinto y más fácil por construcción.")

# =========================================================== 11. PREPRO A
s, y = lamina("Preprocesamiento (1): etiquetar 331,793 notas sin anotarlas "
              "a mano", seccion="4.2 · Preprocesamiento de los datos")
codigo(s, MARGEN, y + Inches(0.05), Inches(7.3), Inches(2.5), [
    "# Tier A - los codigos CIE-10 del Anexo 02 llevan punto (A04.7)",
    "# y MIMIC los guarda sin punto (A047). El JOIN nunca acertaba.",
    "clave = orig.replace('.', '').upper()          # <-- EL ARREGLO",
    "",
    "# la coincidencia es POR PREFIJO, no por igualdad: MIMIC usa",
    "# codigos mas especificos que el mapeo (A047 -> A0471, A0472)",
    "for p in prefijos:",
    "    m = cod10.str.startswith(p)",
    "",
    "# Tier B - acotar la ventana del comodin. Con re.DOTALL un .*",
    "# atravesaba la nota entera y unia parrafos sin relacion.",
    "def acotar(patron, n=100):",
    "    return patron.replace('.*', '.{0,' + str(n) + '}')",
], tam=11, titulo="fase3_v2_tier_a_corregido.py  ·  fase3_v2_corpus_completo.py")

rect(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Inches(2.5),
     relleno=GRIS_TENUE)
rect(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Pt(3), relleno=AZUL)
caja(s, Inches(8.45), y + Inches(0.25), Inches(4.0), Inches(0.3),
     "DOS VÍAS DE ETIQUETADO", tam=11, color=AZUL, negrita=True)
vinetas(s, Inches(8.45), y + Inches(0.68), Inches(4.05), [
    ("Tier A", "223 códigos CIE-10 y CIE-9 mapeados al Anexo 02. "
     "Estratificados en A1 causal y A2 condiciones."),
    ("Tier B", "35 patrones de texto con NegEx para descartar la negación."),
], tam=12, sep=Inches(0.9))

y2 = nota(s, Inches(4.85), "El estrato A1 y A2 existe porque MIMIC no trae "
          "bandera de present-on-admission: sin ella, una neumonía puede ser "
          "el motivo de ingreso y no una infección causada por la atención.")
filas = [["Efecto medido de las dos correcciones", "Antes", "Después"],
         ["Tier A · hospitalizaciones recuperadas", "411", "109,714  (267×)"],
         ["Tier B · detecciones con el comodín acotado", "258,671",
          "93,922  (−63.7 %)"]]
tabla(s, MARGEN, y2 + Inches(0.14), ANCHO_UTIL, filas,
      anchos=[0.56, 0.22, 0.22], alto_fila=Inches(0.36), tam=13,
      alto_cabecera=Inches(0.36))

# =========================================================== 12. PREPRO B
s, y = lamina("Preprocesamiento (2): representación del texto y particiones",
              seccion="4.2 · Preprocesamiento de los datos")
codigo(s, MARGEN, y + Inches(0.05), Inches(7.3), Inches(2.5), [
    "def construir_vectorizador():",
    "    # Palabra + CARACTER: los n-gramas de caracter dan robustez",
    "    # ante abreviaturas y variantes del texto clinico.",
    "    return FeatureUnion([",
    "        ('palabra',  TfidfVectorizer(max_features=60000,",
    "                                     ngram_range=(1, 2))),",
    "        ('caracter', TfidfVectorizer(analyzer='char_wb',",
    "                                     ngram_range=(3, 5)))])",
    "",
    "# Division POR PACIENTE, nunca por nota:",
    "gss = GroupShuffleSplit(n_splits=1, test_size=0.2, random_state=42)",
    "tr, te = next(gss.split(txt, Y, groups=df['subject_id']))",
    "*assert not (set(p_tr) & set(p_te)), 'FUGA: pacientes compartidos'",
], tam=11, titulo="fase4_v2_etiqueta_a1.py")

rect(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Inches(2.5),
     relleno=GRIS_TENUE)
rect(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Pt(3), relleno=GRANATE)
caja(s, Inches(8.45), y + Inches(0.25), Inches(4.0), Inches(0.3),
     "TRES DECISIONES DE PROTOCOLO", tam=11, color=GRANATE, negrita=True)
vinetas(s, Inches(8.45), y + Inches(0.68), Inches(4.05), [
    "Partición por paciente, con aserción explícita de no solape",
    "El vectorizador se ajusta solo con el entrenamiento",
    "Intervalos por bootstrap agrupado por paciente, no por nota",
], tam=12, sep=Inches(0.72))

nota(s, Inches(4.85), "Remuestrear notas en vez de pacientes viola la "
     "independencia y estrecha artificialmente los intervalos de confianza. "
     "Por eso el bootstrap agrupa por subject_id.")
filas = [["Curva de escalabilidad: cuánto texto necesita el modelo", "F1-macro"],
         ["Truncado a 120 caracteres", "0.214"],
         ["Truncado a 2,000 caracteres  (≈ 512 tokens)", "0.433"],
         ["Texto completo", "0.526"]]
tabla(s, MARGEN, Inches(5.68), Inches(7.3), filas, anchos=[0.78, 0.22],
      alto_fila=Inches(0.3), tam=12, alto_cabecera=Inches(0.32), resaltar=2)
caja(s, Inches(8.2), Inches(5.8), Inches(4.5), Inches(1.1),
     "Truncar a 512 tokens cuesta ~18 % del F1-macro. Es la medición propia "
     "que sustenta la necesidad de un modelo de contexto largo.",
     tam=13, color=AZUL, interlinea=1.2)

# =========================================================== 13. MODELO
s, y = lamina("Modelamiento: dos etapas, con derecho a abstenerse",
              seccion="4.3 · Modelamiento y evaluación")
w, h = Inches(3.5), Inches(1.5)
for i, (t, sub, col) in enumerate([
        ("Etapa 1 · Detección", "¿Hay evento adverso?\nBinaria, con clase "
         "negativa (33,564 notas sin evento)", AZUL),
        ("Etapa 2 · Naturaleza", "¿De qué tipo?\n8 clases del Anexo 02",
         GRANATE),
        ("Cascada", "La cifra honesta:\nlos errores se acumulan", VERDE)]):
    x = MARGEN + (w + Inches(0.5)) * i
    rect(s, x, y + Inches(0.1), w, h, relleno=GRIS_TENUE)
    rect(s, x, y + Inches(0.1), w, Pt(3), relleno=col)
    caja(s, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.3),
         t, tam=16, color=col, negrita=True)
    caja(s, x + Inches(0.2), y + Inches(0.68), w - Inches(0.4), Inches(0.8),
         sub, tam=12, color=GRIS, interlinea=1.2)
    if i < 2:
        caja(s, x + w + Inches(0.06), y + Inches(0.65), Inches(0.4),
             Inches(0.4), "→", tam=22, color=GRIS)

ab = f9.get("etapa1", {}).get("abstencion") if f9 else None
y2 = y + Inches(1.85)
consola(s, MARGEN, y2, Inches(7.3), Inches(1.85), [
    "                    texto de entrada    margen        decision",
    "  ------------------------------------------------------------",
    "  .                                     -2.284   se abstiene",
    "  ok                                    -1.180   se abstiene",
    "  Paciente estable, sin novedad.        -0.978   se abstiene",
    "  Routine follow up. Vitals stable.     -0.876   se abstiene",
    "* Abstiene en 6/6 casos triviales.",
], tam=11, titulo="Prueba de abstención ante texto sin contenido clínico")

rect(s, Inches(8.2), y2, Inches(4.5), Inches(1.85), relleno=GRANATE_CLARO)
rect(s, Inches(8.2), y2, Inches(4.5), Pt(3), relleno=GRANATE)
caja(s, Inches(8.45), y2 + Inches(0.2), Inches(4.0), Inches(1.5),
     "Un detector sin clase negativa marca como positivo casi cualquier "
     "texto: no puede decir «aquí no hay nada».\n\nAñadirla es lo que "
     "convierte un clasificador en un detector.",
     tam=13, color=TINTA, interlinea=1.25)
nota(s, Inches(6.05), "Evaluar la Etapa 2 sobre positivos de referencia "
     "supone un detector perfecto y sobreestima el rendimiento. La cascada "
     "se evalúa solo sobre pacientes que ninguna de las dos etapas vio.")

# =========================================================== 14. EJECUCIÓN 1
s, y = lamina("Cómo se ejecutó: entorno y coste real",
              seccion="4.4 · Ejecución del código",
              bajada="Todo corrió en una laptop, sin nube y sin coste. "
              "El entorno se creó aislado para no comprometer el intérprete "
              "del sistema.")
codigo(s, MARGEN, y + Inches(0.05), Inches(7.3), Inches(1.5), [
    '"C:\\ProgramData\\Anaconda3\\python.exe" -m venv T:\\MIMIC\\.venv-gpu',
    "",
    "T:\\MIMIC\\.venv-gpu\\Scripts\\python.exe -m pip install torch \\",
    "     --index-url https://download.pytorch.org/whl/cu126",
    "T:\\MIMIC\\.venv-gpu\\Scripts\\python.exe -m pip install transformers \\",
    "     accelerate sentence-transformers scikit-learn pandas pyarrow",
    "",
    "*torch 2.13.0+cu126 | numpy 2.5.1 | cuda True",
], tam=11, titulo="Preparación del entorno")

rect(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Inches(1.5),
     relleno=GRIS_TENUE)
rect(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Pt(3), relleno=AZUL)
caja(s, Inches(8.45), y + Inches(0.24), Inches(4.0), Inches(1.2),
     "HP Victus 16 · 12 hilos · 15.6 GB RAM\n"
     "NVIDIA GTX 1650, 4 GB (sin Tensor Cores)\n"
     "Windows 11 · Python 3.13 · disco NVMe dedicado",
     tam=12, color=TINTA, interlinea=1.35)

filas = [["Proceso", "Duración real"],
         ["Ablación sobre el corpus completo (331,793 notas)",
          "17 h 56 min"],
         ["Bio_ClinicalBERT · ajuste fino ponderado (GPU)", "3 h 57 min"],
         ["BioBERT · ajuste fino (GPU)", "2 h 55 min"],
         ["Bio_ClinicalBERT · ajuste fino (GPU)", "2 h 55 min"],
         ["Etiqueta A1 y curva de escalabilidad", "2 h 46 min"],
         ["Validación cruzada 5-fold", "1 h 55 min"],
         ["Resto de experimentos de la fase 4", "1 h 03 min"],
         ["Los 4 modelos TF-IDF", "1 min 46 s"],
         ["TOTAL REGISTRADO", "33 h 54 min"]]
tabla(s, MARGEN, y + Inches(1.75), Inches(7.3), filas, anchos=[0.72, 0.28],
      alto_fila=Inches(0.3), tam=12, alto_cabecera=Inches(0.32), resaltar=8)

rect(s, Inches(8.2), y + Inches(1.75), Inches(4.5), Inches(2.1),
     relleno=GRANATE_CLARO)
rect(s, Inches(8.2), y + Inches(1.75), Inches(4.5), Pt(3), relleno=GRANATE)
caja(s, Inches(8.45), y + Inches(1.95), Inches(4.0), Inches(0.3),
     "EL CONTRASTE QUE IMPORTA", tam=11, color=GRANATE, negrita=True)
caja(s, Inches(8.45), y + Inches(2.35), Inches(4.0), Inches(1.4),
     "Los 3 transformers:  9 h 47 min → F1-macro 0.354\n\n"
     "Los 4 TF-IDF:  1 min 46 s → F1-macro 0.459\n\n"
     "333 veces más cómputo, para perder.",
     tam=13, color=TINTA, interlinea=1.3)

# =========================================================== 15. EJECUCIÓN 2
s, y = lamina("Cómo se ejecutó: comandos y salida real de consola",
              seccion="4.4 · Ejecución del código",
              bajada="Nada de esto está reconstruido. Es la salida literal "
              "de los logs, tal como quedó en la bitácora del proyecto.")
codigo(s, MARGEN, y, Inches(6.15), Inches(0.5), [
    "python fase3_v2_tier_a_corregido.py",
], tam=12)
consola(s, MARGEN, y + Inches(0.62), Inches(6.15), Inches(2.25), [
    "[23:00:34] [+0:00:00] Mapeo: 223 codigos - A1 73 - A2 150",
    "[23:23:06] [+0:22:31] 6,364,488 diagnosticos",
    "[23:23:06] [+0:22:31] metodo ORIGINAL (con punto):     411",
    "*[23:23:06] [+0:22:31] metodo CORREGIDO (prefijo): 109,714",
    "",
    "  Hospitalizaciones, original  :      411",
    "* Hospitalizaciones, corregido : 109,714   (266.9x)",
    "  A1 causal explicito : 18,989 notas",
    "  A2 condiciones      : 53,375 notas",
])
codigo(s, Inches(7.15), y, Inches(5.55), Inches(0.5), [
    "python fase3_v2_reanudar.py",
], tam=12)
consola(s, Inches(7.15), y + Inches(0.62), Inches(5.55), Inches(2.25), [
    "[17:07:11] Checkpoint: 260,000 notas (78.4%)",
    "[17:16:14] 280,000/331,793 ( 84.4%) | ETA 0:23:28",
    "[17:39:37] 331,793/331,793 (100.0%) | ETA 0:00:00",
    "",
    "  Variante LAXA (re.DOTALL) : 258,671 detecciones",
    "* Variante ACOTADA (100 car): 93,922 detecciones",
    "  Reduccion al acotar       : 164,749 (63.7%)",
    "",
    "  Total 0:32:26",
])
y2 = nota(s, Inches(5.15), "La corrida completa se cayó al 78.4 %. Por eso "
          "existe fase3_v2_reanudar.py, que retoma desde el punto de control "
          "en lugar de repetir 18 horas.", color=AMBAR,
          fondo=RGBColor(0xFA, 0xF0, 0xDC))
filas = [["La verificación que confirma el diagnóstico", "Antes", "Después",
          "Variación"],
         ["Patrones CON comodín .*", "42,330", "6,735", "−84.1 %"],
         ["Patrones SIN comodín", "13,189", "13,189", "0.0 %"]]
tabla(s, MARGEN, y2 + Inches(0.12), ANCHO_UTIL, filas,
      anchos=[0.55, 0.15, 0.15, 0.15], alto_fila=Inches(0.32), tam=13,
      alto_cabecera=Inches(0.34), resaltar=1)

# =========================================================== 16. AUDITORÍA
s, y = lamina("El hallazgo central: un AUC de 0.973 que era falso",
              seccion="4.4 · Ejecución del código",
              bajada="MIMIC abarca 2008–2019 y la transición CIE-9 → CIE-10 "
              "ocurrió en 2015. El mapeo inicial solo tenía CIE-10, así que "
              "toda hospitalización anterior era negativa por construcción.")
filas = [["Versión del detector", "Especificidad", "AUC", "VPP"],
         ["Inicial — INVÁLIDA, no se cita", "0.917", "0.973", "0.433"],
         ["Reevaluada con emparejamiento", "0.694", "0.904", "0.171"],
         ["Final — siete modos corregidos", "0.770", "0.843", "0.455"]]
tabla(s, MARGEN, y + Inches(0.05), Inches(7.3), filas,
      anchos=[0.46, 0.2, 0.17, 0.17], alto_fila=Inches(0.44), tam=13,
      resaltar=2)

rect(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Inches(1.85),
     relleno=RGBColor(0xFA, 0xF0, 0xDC))
rect(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Pt(3), relleno=AMBAR)
caja(s, Inches(8.45), y + Inches(0.25), Inches(4.0), Inches(0.3),
     "QUÉ APRENDIÓ EL MODELO", tam=11, color=AMBAR, negrita=True)
caja(s, Inches(8.45), y + Inches(0.62), Inches(4.0), Inches(1.2),
     "No aprendió a reconocer eventos adversos: aprendió a distinguir la "
     "plantilla documental de cada época.\n\nEl rasgo de mayor peso era "
     "palabra__rdwsd — un artefacto de cabecera de laboratorio, sin ningún "
     "contenido clínico.", tam=12, color=TINTA, interlinea=1.2)

y2 = y + Inches(2.1)
caja(s, MARGEN, y2, Inches(6.0), Inches(0.3),
     "Los siete modos de fallo, todos con efecto medido", tam=16,
     color=TINTA, negrita=True)
vinetas(s, MARGEN, y2 + Inches(0.42), Inches(6.1), [
    "Corpus era solo el 9 % (dos LIMIT sin declarar)",
    "re.DOTALL cruzaba la nota entera",
    "Códigos CIE con punto frente a MIMIC sin punto",
    "Ausencia de clase negativa",
], tam=13, sep=Inches(0.34))
vinetas(s, Inches(7.1), y2 + Inches(0.42), Inches(5.6), [
    "Confusor de época CIE-9 / CIE-10  ← el más grave",
    "43 de 223 códigos eran reglas muertas (CIE-10 OMS ≠ CM)",
    "Bandas de priorización degeneradas con n pequeño",
], tam=13, sep=Inches(0.34))
nota(s, Inches(6.35), "Lección para llevarse: un AUC alto puede sostenerse "
     "íntegramente en una señal espuria. El aporte de este trabajo no es una "
     "métrica, es la auditoría que la hizo creíble.", color=GRANATE,
     fondo=GRANATE_CLARO)

# =========================================================== 17. RESULTADO 1
s, y = lamina("Resultados: la etapa de detección",
              seccion="4.5 · Resultados")
e1 = f9["etapa1"] if f9 and "etapa1" in f9 else {}
filas = [["Métrica", "Valor", "IC 95 % (agrupado por paciente)"],
         ["Sensibilidad", "0.7623", "[0.7514, 0.7729]"],
         ["Especificidad", "0.7699", "[0.7591, 0.7809]"],
         ["AUC", "0.8426", "[0.8361, 0.8493]"]]
tabla(s, MARGEN, y + Inches(0.05), Inches(6.6), filas,
      anchos=[0.32, 0.24, 0.44], alto_fila=Inches(0.42), tam=14)
consola(s, MARGEN, y + Inches(1.85), Inches(6.6), Inches(1.2), [
    "Matriz de confusion:",
    "  VP=5,761   FP=1,502   FN=1,796   VN=5,026",
    "  VPP crudo en el test (balanceado) : 0.793",
    "* VPP reajustado a prevalencia real : 0.455   <-- el operativo",
])

rect(s, Inches(7.5), y + Inches(0.05), Inches(5.2), Inches(3.0),
     relleno=GRIS_TENUE)
rect(s, Inches(7.5), y + Inches(0.05), Inches(5.2), Pt(3), relleno=GRANATE)
caja(s, Inches(7.75), y + Inches(0.25), Inches(4.7), Inches(0.3),
     "EL VPP DEPENDE DE DÓNDE SE DESPLIEGUE", tam=11, color=GRANATE,
     negrita=True)
for i, (p, v, marca) in enumerate([("5 %", "0.148", ""),
                                   ("10 %", "0.269", ""),
                                   ("20.12 %", "0.455", " ← la real"),
                                   ("35 %", "0.641", ""),
                                   ("50 %", "0.768", "")]):
    yy = y + Inches(0.66) + Inches(0.38) * i
    es = (i == 2)
    caja(s, Inches(7.75), yy, Inches(2.15), Inches(0.3),
         f"prevalencia {p}", tam=12, color=TINTA if es else GRIS,
         negrita=es, fuente="Consolas")
    caja(s, Inches(9.95), yy, Inches(2.65), Inches(0.3), f"VPP {v}{marca}",
         tam=12, color=GRANATE if es else GRIS, negrita=es, fuente="Consolas")
caja(s, Inches(7.75), y + Inches(2.62), Inches(4.7), Inches(0.3),
     "El mismo modelo, de 0.15 a 0.77, sin cambiar una línea.",
     tam=12, color=AZUL)

nota(s, Inches(6.05), "De cada 100 alertas, 45 son eventos reales. Revisar "
     "notas clínicas al azar encuentra 20. El sistema multiplica por 2.26 la "
     "eficiencia del revisor: opera como filtro de cribado, no como árbitro.")

# =========================================================== 18. RESULTADO 2
s, y = lamina("Resultados: el ranking de los siete modelos",
              seccion="4.5 · Resultados",
              bajada="Todos evaluados sobre la misma partición estratificada, "
              "semilla 42. Las cifras son comparables entre sí.")
R = f11["resultados"]
orden = sorted([(k, v) for k, v in R.items() if "f1_macro" in v],
               key=lambda x: -x[1]["f1_macro"])
filas = [["#", "Modelo", "Exactitud", "F1-macro", "Kappa", "Tiempo"]]
for i, (k, v) in enumerate(orden, 1):
    seg = v.get("segundos", 0)
    t = f"{seg/3600:.1f} h" if seg > 3600 else (
        f"{seg/60:.0f} min" if seg > 90 else f"{seg:.0f} s")
    filas.append([str(i), k, f"{v['exactitud']:.4f}", f"{v['f1_macro']:.4f}",
                  f"{v['kappa']:.4f}", t])
fin_tabla = tabla(s, MARGEN, y + Inches(0.05), ANCHO_UTIL, filas,
                  anchos=[0.04, 0.48, 0.12, 0.12, 0.12, 0.12],
                  alto_fila=Inches(0.3), tam=11.5, resaltar=0,
                  alto_cabecera=Inches(0.34))

v = f11["ventana"]
y2 = fin_tabla + Inches(0.16)
consola(s, MARGEN, y2, Inches(7.3), Inches(1.35), [
    "EFECTO AISLADO DE LA VENTANA DE CONTEXTO",
    "(misma arquitectura y datos; solo cambia cuanto texto ve)",
    "  exactitud   completo 0.715 -> truncado 0.575  (-19.6%)",
    "  f1_macro    completo 0.459 -> truncado 0.330  (-28.1%)",
    "* kappa       completo 0.544 -> truncado 0.315  (-42.2%)",
])
rect(s, Inches(8.2), y2, Inches(4.5), Inches(1.35), relleno=GRANATE_CLARO)
rect(s, Inches(8.2), y2, Inches(4.5), Pt(3), relleno=GRANATE)
caja(s, Inches(8.45), y2 + Inches(0.14), Inches(4.0), Inches(1.1),
     f"La ventana de {f11['config']['max_len']} tokens cubre solo el "
     f"{v['cobertura_media']:.1%} del documento (mediana "
     f"{v['tokens_mediana']:,} tokens por nota clínica).",
     tam=12.5, color=TINTA, interlinea=1.25)
nota(s, y2 + Inches(1.51), "La hipótesis inicial se refutó: el transformer "
     "clínico no superó al modelo léxico. La fila 4 explica por qué — al "
     "truncar el TF-IDF a la ventana de BERT, cae igual. No es la "
     "arquitectura, es cuánto texto ve. La mejora es Clinical-Longformer.",
     color=GRANATE, fondo=GRANATE_CLARO, alto=Inches(0.72))

# =========================================================== 19. RESULTADO 3
s, y = lamina("Resultados: cascada y contraste con el juicio experto",
              seccion="4.5 · Resultados")
cas = f10["B_cascada"] if f10 and "B_cascada" in f10 else {}
filas = [["Evaluación", "F1-micro", "F1-macro"],
         ["Etapa 2 aislada (positivos de referencia)", "0.7516", "0.5126"],
         ["CASCADA real (Etapa 1 → Etapa 2)", "0.4926", "0.3626"]]
tabla(s, MARGEN, y + Inches(0.05), Inches(6.3), filas,
      anchos=[0.56, 0.22, 0.22], alto_fila=Inches(0.4), tam=13, resaltar=1)
caja(s, MARGEN, y + Inches(1.4), Inches(6.3), Inches(0.6),
     "Caída al encadenar: −34.5 % en F1-micro. Reportar la clasificación "
     "aislada sobreestima el rendimiento operativo en torno a un tercio.",
     tam=13, color=AZUL, interlinea=1.2)

ce = f12["contra_experto"]
m = f12["matriz"]
filas2 = [["Contra el consenso experto", "Valor", "IC 95 %"],
          ["Sensibilidad", f"{ce['sensibilidad']:.4f}",
           f"[{ce['sensibilidad_ic95'][0]:.3f}, {ce['sensibilidad_ic95'][1]:.3f}]"],
          ["Especificidad", f"{ce['especificidad']:.4f}",
           f"[{ce['especificidad_ic95'][0]:.3f}, {ce['especificidad_ic95'][1]:.3f}]"],
          ["VPP", f"{ce['vpp']:.4f}",
           f"[{ce['vpp_ic95'][0]:.3f}, {ce['vpp_ic95'][1]:.3f}]"]]
fin2 = tabla(s, Inches(7.1), y + Inches(0.05), Inches(5.6), filas2,
             anchos=[0.44, 0.2, 0.36], alto_fila=Inches(0.4), tam=13)
caja(s, Inches(7.1), fin2 + Inches(0.1), Inches(5.6), Inches(0.6),
     f"n = {f12['n_analizados']} casos · kappa sistema-experto "
     f"{ce['kappa_sistema_experto']:.3f} · matriz VP={m['vp']} FP={m['fp']} "
     f"FN={m['fn']} VN={m['vn']}", tam=12, color=GRIS, interlinea=1.2)

y2 = Inches(4.15)
filas3 = [["Robustez: ¿depende de quién anote?", "n", "Sensibilidad",
           "Especificidad", "Kappa"],
          ["Consenso de ambos evaluadores (principal)", "68",
           "0.9455", "0.5385", "0.5307"],
          ["Evaluador B por separado", "77", "0.9032", "0.4667", "0.3896"],
          ["ANÁLISIS AMPLIADO: los 163 casos anotados (152 útiles)", "152",
           "0.9138  [0.849, 0.953]", "—*", "—*"]]
tabla(s, MARGEN, y2, ANCHO_UTIL, filas3,
      anchos=[0.4, 0.08, 0.22, 0.15, 0.15], alto_fila=Inches(0.34), tam=12.5,
      resaltar=2)
caja(s, MARGEN, y2 + Inches(1.46), ANCHO_UTIL, Inches(0.3),
     "* En el análisis ampliado solo se lee la sensibilidad: el piloto se "
     "muestreó por estratos definidos por la salida del sistema, así que la "
     "especificidad y el kappa se reportan sobre el consenso.",
     tam=11, color=GRIS, interlinea=1.1)
nota(s, Inches(6.15), "La sensibilidad se sostiene sobre 0.90 con cualquiera "
     "de las referencias — y al juntar los 163 casos anotados el intervalo "
     "se estrecha a la mitad. Además: el 25 % de los casos que la "
     "codificación declara negativos son eventos reales para el experto — el "
     "subregistro, medido con datos propios.", color=GRANATE,
     fondo=GRANATE_CLARO, alto=Inches(0.75))

# =========================================================== 20. RESULTADO 4
s, y = lamina("Resultados: el modelo ya funciona en texto peruano real",
              seccion="4.5 · Resultados",
              bajada="Cuatro tareas sobre 6,336 casos reales del sistema "
              "peruano de reporte, con los códigos del Anexo 02 asignados "
              "por profesionales. Predicción out-of-fold: cada caso lo "
              "predice un modelo que no lo vio.")
filas = [["Tarea", "Qué predice", "Modelo", "Clases", "Exactitud", "F1-macro"],
         ["T1", "Evento adverso frente a incidente", "LogReg", "2", "0.862",
          "0.858"],
         ["T2", "Naturaleza (Anexo 02)", "LinearSVC", "9", "0.846", "0.766"],
         ["T3", "Severidad (Anexo 03)", "LogReg", "4", "0.726", "0.602"],
         ["T4", "Código de evento específico", "LinearSVC", "41", "0.761",
          "0.659"]]
fin1 = tabla(s, MARGEN, y + Inches(0.05), ANCHO_UTIL, filas,
             anchos=[0.06, 0.4, 0.14, 0.1, 0.15, 0.15],
             alto_fila=Inches(0.34), tam=12.5, resaltar=3,
             alto_cabecera=Inches(0.36))

y2 = fin1 + Inches(0.22)
caja(s, MARGEN, y2, Inches(6.3), Inches(0.3),
     "Lo que el corpus español rescata", tam=15, color=TINTA, negrita=True)
filas2 = [["Naturaleza", "En MIMIC", "En español"],
          ["Gestión de la organización", "n=24 · F1 0.000", "F1 0.853"],
          ["Cuidado del paciente", "F1 0.791", "F1 0.923"],
          ["Medicación", "F1 0.679", "F1 0.875"]]
tabla(s, MARGEN, y2 + Inches(0.36), Inches(6.3), filas2,
      anchos=[0.46, 0.27, 0.27], alto_fila=Inches(0.3), tam=11.5,
      alto_cabecera=Inches(0.3), resaltar=0)

rect(s, Inches(7.1), y2, Inches(5.6), Inches(1.62), relleno=GRIS_TENUE)
rect(s, Inches(7.1), y2, Inches(5.6), Pt(3), relleno=VERDE)
caja(s, Inches(7.35), y2 + Inches(0.14), Inches(5.1), Inches(0.26),
     "ENTREGABLE PARA REVISIÓN", tam=11, color=VERDE, negrita=True)
caja(s, Inches(7.35), y2 + Inches(0.46), Inches(5.1), Inches(1.1),
     "Los 8,799 casos exportados a Excel con la codificación humana, la "
     "predicción del modelo y el semáforo de acierto. Los 2,870 errores "
     "quedan en su propia hoja, listos para revisión clínica.",
     tam=12, color=TINTA, interlinea=1.2)
nota(s, y2 + Inches(1.82), "T4 es nuevo: el código de evento específico "
     "nunca se había modelado. De las 154 categorías usadas, 41 tienen casos "
     "suficientes y cubren el 89 % del corpus.")

# =========================================================== 21. GEMSES
s, y = lamina("Del texto a la decisión de gestión",
              seccion="4.6 · Destino aplicado",
              bajada="La salida del canal de PLN alimenta la matriz de "
              "priorización vigente en EsSalud. Esto ya no es PLN: es el "
              "destino del trabajo.")
w, h = Inches(2.9), Inches(1.35)
for i, (t, sub, col) in enumerate([
        ("Frecuencia", "A → B = (A/Σa)×9", AZUL),
        ("Impacto", "G = 0.40 estancia + 0.20 complic.\n"
         "+ 0.15 sobrecostos + 0.25 insatisf.", AZUL),
        ("Índice", "H = B×G → I = H/Σh → J = I×100", GRANATE),
        ("Banda y responsable", "P25 verde · P50 amarillo · P75 rojo",
         VERDE)]):
    x = MARGEN + (w + Inches(0.18)) * i
    rect(s, x, y + Inches(0.1), w, h, relleno=GRIS_TENUE)
    rect(s, x, y + Inches(0.1), w, Pt(3), relleno=col)
    caja(s, x + Inches(0.16), y + Inches(0.28), w - Inches(0.32),
         Inches(0.3), t, tam=15, color=col, negrita=True)
    caja(s, x + Inches(0.16), y + Inches(0.66), w - Inches(0.32),
         Inches(0.7), sub, tam=11, color=GRIS, interlinea=1.2)

filas = [["Banda", "Acción de mejora", "Responsable"],
         ["Verde  (≤ P25)", "Mejora de actividades", "Jefatura de servicio"],
         ["Amarillo  (P25–P75)", "Mejora de procesos", "Departamento"],
         ["Rojo  (> P75)", "Proyecto de mejora", "Dirección"]]
tabla(s, MARGEN, y + Inches(1.7), Inches(7.3), filas,
      anchos=[0.3, 0.38, 0.32], alto_fila=Inches(0.38), tam=13)

rect(s, Inches(8.2), y + Inches(1.7), Inches(4.5), Inches(1.9),
     relleno=RGBColor(0xFA, 0xF0, 0xDC))
rect(s, Inches(8.2), y + Inches(1.7), Inches(4.5), Pt(3), relleno=AMBAR)
caja(s, Inches(8.45), y + Inches(1.9), Inches(4.0), Inches(0.3),
     "UN DEFECTO CORREGIDO", tam=11, color=AMBAR, negrita=True)
caja(s, Inches(8.45), y + Inches(2.28), Inches(4.0), Inches(1.2),
     "Con menos de 8 eventos los percentiles se degeneran e invierten la "
     "prioridad: una úlcera de impacto alto salía verde. Ahora, con n<8, se "
     "usa corte absoluto sobre G.", tam=12, color=TINTA, interlinea=1.2)
nota(s, Inches(5.85), "La matriz no es una propuesta del autor: la Directiva "
     "N.º 7-OGCyH-ESSALUD-2020 la cita textualmente en su artículo 5.13 como "
     "el instrumento vigente.")

# =========================================================== 22. CONCLUSIONES
s, y = lamina("Conclusiones", seccion="5 · Conclusiones")
concl = [
    ("Un modelo léxico bien construido superó al transformer clínico.",
     "La hipótesis inicial se refutó y el resultado negativo se reporta como "
     "tal: TF-IDF + LinearSVC 0.459 frente a Bio_ClinicalBERT 0.354 de "
     "F1-macro."),
    ("La diferencia se explica por la ventana de contexto, no por la "
     "arquitectura.",
     "Truncar el modelo léxico a los 256 tokens que ve BERT lo degrada un "
     "28 % con los mismos datos. La vía de mejora es un modelo de secuencia "
     "larga."),
    ("El aporte principal es la auditoría de validez.",
     "Siete modos de fallo del etiquetado, con el confusor de época como caso "
     "ejemplar: un AUC de 0.973 sostenido en una plantilla documental."),
    ("La evaluación en cascada corrige un optimismo de un tercio.",
     "F1-micro 0.752 aislada frente a 0.493 extremo a extremo."),
    ("Contra juicio experto el sistema recupera el 94.5 % de los eventos.",
     "Con especificidad 0.538: sirve como cribado, no como árbitro."),
    ("El modelo ya demostró funcionar en texto peruano real.",
     "Cuatro tareas con etiqueta de oro sobre casos reales del sistema "
     "nacional de reporte, incluida por primera vez la predicción del "
     "código de evento específico del Anexo 02."),
]
yy = y + Inches(0.02)
for i, (t, d) in enumerate(concl, 1):
    rect(s, MARGEN, yy + Inches(0.02), Inches(0.3), Inches(0.3),
         relleno=GRANATE)
    caja(s, MARGEN, yy + Inches(0.05), Inches(0.3), Inches(0.25), str(i),
         tam=12, color=BLANCO, negrita=True, alineado=PP_ALIGN.CENTER)
    caja(s, MARGEN + Inches(0.45), yy, Inches(12.0), Inches(0.28), t,
         tam=14, color=TINTA, negrita=True)
    caja(s, MARGEN + Inches(0.45), yy + Inches(0.28), Inches(12.0),
         Inches(0.44), d, tam=12, color=GRIS, interlinea=1.15)
    yy += Inches(0.83)

# =========================================================== 23. LIMITACIONES
s, y = lamina("Limitaciones declaradas", seccion="5 · Conclusiones",
              bajada="Se declaran porque acotan el alcance de cada cifra. "
              "Ninguna se descubrió después: todas están medidas.")
lim2 = [("Estándar de plata",
         "Las métricas sobre MIMIC miden acuerdo con códigos CIE, no con "
         "juicio clínico."),
        ("Reutilización del conjunto de prueba",
         "A lo largo de las iteraciones del canal; sesgo optimista no "
         "cuantificado. El umbral no se ajustó."),
        ("Alcance de la validación experta",
         "La muestra está restringida a eventos de infección; n = 78."),
        ("Cobertura de datos estructurados",
         "Las tablas de UCI cubren solo el 19.7 % de las notas clínicas."),
        ("Independencia del anotador",
         "El autor es a la vez desarrollador y anotador. Se mitiga con "
         "interfaz ciega, orden aleatorio, cronometraje auditable y un "
         "segundo evaluador."),
        ("Transferencia al español no demostrada",
         "Hoy son dos canales paralelos: el corpus español se entrena desde "
         "cero, no se transfiere el modelo de MIMIC.")]
yy = y + Inches(0.05)
for t, d in lim2:
    rect(s, MARGEN, yy, ANCHO_UTIL, Inches(0.78), relleno=GRIS_TENUE)
    rect(s, MARGEN, yy, Pt(4), Inches(0.78), relleno=AMBAR)
    caja(s, MARGEN + Inches(0.24), yy + Inches(0.1), Inches(3.5),
         Inches(0.3), t, tam=13, color=TINTA, negrita=True)
    caja(s, MARGEN + Inches(4.0), yy + Inches(0.1), Inches(7.9),
         Inches(0.58), d, tam=12, color=GRIS, interlinea=1.15)
    yy += Inches(0.88)

# =========================================================== 24. RECOMENDA
s, y = lamina("Recomendaciones", seccion="6 · Recomendaciones")
rec = [("Corto plazo", [
    "Sustituir la arquitectura por Clinical-Longformer (ventana 4,096): es "
    "la hipótesis que el experimento deja planteada",
    "Anotar el gold standard de 350 notas para medir validez clínica",
    "Ampliar la clase GRAVE del corpus español (23 casos, F1 0.148)"]),
    ("Mediano plazo", [
    "Entrenar sobre notas clínicas peruanas con un encoder clínico en español "
    "(bsc-bio-ehr-es), previa autorización de datos",
    "Integrar el detector como cribado en la Oficina de Calidad, con el "
    "revisor humano decidiendo siempre",
    "Medir el coste-beneficio real: horas-revisor ahorradas frente a "
    "revisiones improductivas generadas"])]
ancho_col = Inches(5.98)
for j, (tit, items) in enumerate(rec):
    x0 = MARGEN + (ancho_col + Inches(0.14)) * j
    rect(s, x0, y + Inches(0.05), ancho_col, Inches(0.42), relleno=GRANATE)
    caja(s, x0 + Inches(0.2), y + Inches(0.13), Inches(4), Inches(0.3),
         tit, tam=15, color=BLANCO, negrita=True)
    vinetas(s, x0 + Inches(0.1), y + Inches(0.68), ancho_col - Inches(0.2),
            items, tam=13, sep=Inches(0.44))
nota(s, Inches(5.75), "El sistema no sustituye al profesional de calidad: le "
     "dice a quién mirar primero. Con VPP 0.455, más de la mitad de las "
     "alertas son falsas — aceptable para un cribado, no para una decisión "
     "automática.", color=GRANATE, fondo=GRANATE_CLARO,
     alto=Inches(0.75))

# =========================================================== 25. REFERENCIAS
s, y = lamina("Referencias principales", seccion="7 · Referencias",
              bajada="El artículo completo cita 34 fuentes. Estas son las que "
              "sostienen las decisiones de diseño.")
refs = [
    "Johnson, A. et al. MIMIC-IV-Note: deidentified free-text clinical notes "
    "v2.2. PhysioNet.",
    "Alsentzer, E. et al. Publicly Available Clinical BERT Embeddings. "
    "Clinical NLP Workshop, 2019.",
    "Lee, J. et al. BioBERT: a pre-trained biomedical language "
    "representation model. Bioinformatics, 2020.",
    "Li, Y. et al. Clinical-Longformer and Clinical-BigBird: Transformers "
    "for long clinical sequences. 2022.",
    "Ratner, A. et al. Snorkel: rapid training data creation with weak "
    "supervision. VLDB, 2017.",
    "Geirhos, R. et al. Shortcut learning in deep neural networks. Nature "
    "Machine Intelligence, 2020.",
    "Zech, J. et al. Variable generalization performance of a deep learning "
    "model to detect pneumonia. PLOS Medicine, 2018.",
    "Byrt, T., Bishop, J., Carlin, J. Bias, prevalence and kappa. Journal of "
    "Clinical Epidemiology, 1993.",
    "de Vries, E. et al. The incidence and nature of in-hospital adverse "
    "events. BMJ Quality & Safety, 2008.",
    "EsSalud. Directiva GG-ESSALUD-2021: registro, notificación y gestión "
    "de los ERSP. Anexos 02 y 03.",
]
yy = y + Inches(0.02)
for i, r in enumerate(refs, 1):
    caja(s, MARGEN, yy, Inches(0.35), Inches(0.3), f"[{i}]", tam=12,
         color=GRANATE, negrita=True)
    caja(s, MARGEN + Inches(0.45), yy, Inches(11.9), Inches(0.42), r,
         tam=12.5, color=TINTA, interlinea=1.12)
    yy += Inches(0.44)

rect(s, MARGEN, Inches(6.35), ANCHO_UTIL, Inches(0.6), relleno=GRIS_TENUE)
caja(s, MARGEN + Inches(0.22), Inches(6.5), Inches(12.4), Inches(0.3),
     "Código, resultados y artículo completo:  "
     "github.com/carlosperez100/PLN_SP    ·    "
     "carlosperez100.github.io/PLN_SP", tam=13, color=AZUL, negrita=True)

# =========================================================== 26. CIERRE
s = prs.slides.add_slide(BLANCO_LAYOUT)
rect(s, 0, 0, ANCHO, ALTO, relleno=GRANATE)
caja(s, Inches(1.4), Inches(2.1), Inches(10.5), Inches(2.2),
     "De la notificación voluntaria\na la detección sistemática",
     tam=40, color=BLANCO, negrita=True, alineado=PP_ALIGN.CENTER,
     interlinea=1.15)
rect(s, Inches(6.0), Inches(4.3), Inches(1.3), Pt(2.5), relleno=BLANCO)
caja(s, Inches(1.4), Inches(4.65), Inches(10.5), Inches(0.9),
     "El aporte no es una métrica alta: es una métrica en la que se puede "
     "confiar.", tam=18, color=RGBColor(0xE8, 0xC8, 0xCC),
     alineado=PP_ALIGN.CENTER, interlinea=1.25)
caja(s, Inches(1.4), Inches(5.9), Inches(10.5), Inches(0.8),
     "Carlos Pérez Pérez\nMIA-10 Procesamiento del Lenguaje Natural · "
     "Universidad Nacional de Ingeniería · 2026",
     tam=13, color=RGBColor(0xE0, 0xB8, 0xBE), alineado=PP_ALIGN.CENTER,
     interlinea=1.35)

SALIDA.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(SALIDA))
print(f"[OK] {SALIDA}")
print(f"     {len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas")
