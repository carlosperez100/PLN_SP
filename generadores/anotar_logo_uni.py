# -*- coding: utf-8 -*-
"""Estampa el logo de la UNI en todas las laminas del PPTX final.

Se ejecuta DESPUES de generar_presentacion_final.py y anotar_notas_orador.py:
    python anotar_logo_uni.py [ruta_pptx_opcional]

Posicion: esquina superior derecha, discreto (0.62" de alto), en las 29
laminas. En la portada va mas grande (1.15") porque alli es protagonista.
"""
import io
import sys
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8",
                              errors="replace")
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

LOGO = r"T:\MIMIC\PLN_SP\paper\figuras\logo_uni.png"
P = sys.argv[1] if len(sys.argv) > 1 else (
    r"T:\MIMIC\PLN_SP\presentacion"
    r"\Presentacion_Trabajo_Final_PLN_CarlosPerez.pptx")

ANCHO_LAMINA = Inches(13.333)
im = Image.open(LOGO)
aspecto = im.size[0] / im.size[1]          # ancho / alto

pr = Presentation(P)
for i, s in enumerate(pr.slides):
    if i == 0:                             # portada: mas grande
        h = Inches(1.15)
        w = Inches(1.15 * aspecto)
        x = ANCHO_LAMINA - w - Inches(0.45)
        y = Inches(0.35)
    else:                                  # resto: discreto arriba-derecha
        h = Inches(0.62)
        w = Inches(0.62 * aspecto)
        x = ANCHO_LAMINA - w - Inches(0.3)
        y = Inches(0.22)
    s.shapes.add_picture(LOGO, x, y, width=w, height=h)

pr.save(P)
print(f"[OK] logo UNI estampado en {len(pr.slides)} laminas -> {P}")
