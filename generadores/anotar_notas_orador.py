# -*- coding: utf-8 -*-
"""Incrusta el guion hablado como NOTAS DEL ORADOR en cada lamina del PPTX.

Las notas se ven en la vista Moderador de PowerPoint (Alt+F5) y al imprimir
"Paginas de notas"; el publico nunca las ve. Se ejecuta DESPUES del generador:

    python generar_presentacion_final.py
    python anotar_notas_orador.py
"""
import io
import sys
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8",
                              errors="replace")
from pptx import Presentation

# por defecto anota el master; acepta otra ruta como primer argumento
P = sys.argv[1] if len(sys.argv) > 1 else (
    r"T:\MIMIC\PLN_SP\presentacion"
    r"\Presentacion_Trabajo_Final_PLN_CarlosPerez.pptx")

NOTAS = [
    # 1 · Portada
    "Buenas tardes. Mi trabajo se titula Deteccion automatica de eventos "
    "adversos hospitalarios en notas clinicas mediante PLN. Trabajo final "
    "del curso MIA-10 con el Dr. Wester Zela. Cuatro cifras que resumen lo "
    "que van a ver: 70,000 notas clinicas de modelado, casi 34 horas de "
    "computo medido, 7 modelos comparados, y un AUC final de 0.843 — parte "
    "de la historia es por que ese numero es MENOR que uno que obtuve "
    "antes, y por que eso es una mejora. (30 segundos, no mas.)",

    # 2 · Agenda
    "Siete estaciones: contexto, planteamiento del problema, antecedentes, "
    "y la parte central — la solucion — con datos, preprocesamiento, "
    "modelamiento, COMO SE EJECUTO EL CODIGO con sus costos reales, y "
    "resultados. Cierro con conclusiones, recomendaciones y referencias. "
    "(Subrayar con la voz 'como se ejecuto el codigo': es lo que pidio el "
    "docente.)",

    # 3 · Introduccion
    "El concepto primero: un evento adverso es un dano NO intencional "
    "causado por la atencion de salud, no por la enfermedad. Ejemplo: una "
    "neumonia puede ser el motivo de ingreso (no es evento) o adquirida "
    "por el ventilador durante la estancia (si lo es). Esa distincion es "
    "dificil hasta para un humano y persigue todo el trabajo. La vara de "
    "medir es normativa: el Anexo 02 de la Directiva GG-ESSALUD-2021, con "
    "231 eventos en 12 naturalezas. No invente categorias: uso las "
    "oficiales.",

    # 4 · Problematica
    "La notificacion es manual y voluntaria; el subregistro estimado es "
    "cercano al 72%: de cada 10 eventos, 7 nunca llegan a calidad. La "
    "paradoja que motiva el trabajo: la informacion SI existe — el medico "
    "que atendio la complicacion la escribio en la nota clinica. Lo que no "
    "existe es la capacidad humana de leer el 100% de las notas. Pregunta "
    "de investigacion: puede un sistema de PLN leerlas y detectar los "
    "eventos con fiabilidad comparable al juicio experto? (Hablar desde la "
    "experiencia de gestion: es mi terreno.)",

    # 5 · Objetivos
    "Cinco objetivos especificos, todos cumplidos: corpus por supervision "
    "debil (70,000 notas); comparacion clasico vs transformers (7 "
    "modelos); priorizacion institucional; validacion experta (163 casos "
    "con juicio experto); y transferencia al espanol (hecha, con datos "
    "peruanos reales). Transparencia: la hipotesis central — que el "
    "transformer superaria al modelo lexico — SE REFUTO. Un resultado "
    "negativo bien medido es un resultado: mas adelante muestro por que "
    "perdio y cuanto costo comprobarlo. Si preguntan si es un fracaso: lo "
    "seria no poder explicar por que fallo; medi la causa (la ventana de "
    "contexto) y eso define la linea de mejora.",

    # 6 · Antecedentes
    "Tres cuerpos de literatura: (1) Supervision debil — paradigma Snorkel "
    "de Stanford: etiquetar con reglas programaticas acepta ruido a cambio "
    "de escala; la validez de las reglas HAY que auditarla — corazon de mi "
    "trabajo. (2) Aprendizaje por atajo — Geirhos 2020: el modelo resuelve "
    "por la via estadisticamente barata, no la correcta. Caso Zech 2018: "
    "detector de neumonia con AUC altisimo que reconocia DE QUE HOSPITAL "
    "venia la radiografia. Retengan ese caso: me paso lo mismo y lo voy a "
    "mostrar. (3) Ventana de contexto: Clinical-Longformer supera a "
    "ClinicalBERT al ampliar la ventana — justo lo que mis datos "
    "confirman por otra via.",

    # 7 · Seleccion del modelo
    "Antes de comparar hay que decidir que modelos son ADMISIBLES: filtros "
    "duros documentados. El F4 es el que casi nadie aplica: PhysioNet "
    "prohibe expresamente enviar texto de MIMIC a APIs de terceros "
    "(ChatGPT incluido) — condicion legal firmada, no preferencia. Eso "
    "descarta los LLM comerciales de entrada. F3: un LLM local de 7B no "
    "cabe en el hardware de despliegue. Dato verificable: el Open LLM "
    "Leaderboard esta ARCHIVADO por sus propios autores; citarlo hoy seria "
    "un error comprobable. Use el ranking pertinente: clasificacion, no "
    "generacion. La tarea determina la metrica y la metrica el ranking — "
    "no al reves.",

    # 8 · Flujo completo (imagen)
    "Este diagrama es el mapa de todo; recorrerlo banda por banda con el "
    "puntero (90 seg). BANDA 1 — corpus: 331,793 notas de MIMIC-IV mas sus "
    "codigos; supervision debil con dos vias: Tier A cruza codigos CIE "
    "contra el Anexo 02, Tier B busca 35 patrones de texto con filtro de "
    "negacion NegEx. La caja amarilla: la auditoria que corrigio 7 modos "
    "de fallo. Sale el corpus de 70,000 notas, prevalencia real 20.12%. "
    "BANDA 2 — modelado: TF-IDF palabra+caracter, particion POR PACIENTE; "
    "dos etapas (detectar con abstencion, clasificar naturaleza); "
    "evaluacion honesta: cascada extremo a extremo y juicio experto en 163 "
    "casos. BANDA 3 — transferencia: 8,799 reportes peruanos reales con "
    "etiqueta de oro; destino: matriz de priorizacion institucional. "
    "Cerrar: 'todo lo que sigue es entrar a cada caja de este mapa'.",

    # 9 · Fases y scripts
    "Cada caja del diagrama tiene su script con nombre y apellido; todo el "
    "codigo es publico en github.com/carlosperez100/PLN_SP. Lo que NUNCA "
    "se versiona es el dato clinico (texto, modelos entrenados sobre "
    "texto): DUA de PhysioNet y proteccion de datos. La bitacora guarda "
    "comando y salida de consola de cada corrida. NO leer la tabla — "
    "queda para preguntas. (15 segundos y avanzar.)",

    # 10 · Datos MIMIC
    "MIMIC-IV-Note: 331,793 resumenes de alta reales del Beth Israel "
    "Deaconess de Boston, 2008-2019, acceso credencializado (curso CITI + "
    "DUA). Texto clinico real con toda su suciedad. El corpus de modelado "
    "quedo en 70,000 notas: el techo lo puso la RAM al vectorizar "
    "n-gramas de caracter, y esta declarado. Numero clave de la lamina: "
    "prevalencia poblacional real 20.12% (109,775 hospitalizaciones con "
    "evento sobre 545,601). Retenerlo: decide el VPP operativo. Si "
    "preguntan por que datos americanos: no existe corpus peruano de "
    "notas con eventos etiquetados y acceso autorizado; MIMIC es el "
    "estandar mundial para desarrollar el metodo legalmente — y la "
    "validacion en espanol viene en dos laminas.",

    # 11 · ERSP
    "La otra mitad, y es peruana: 8,799 ocurrencias REALES del sistema "
    "institucional de reporte, en espanol, codificadas UNA A UNA por "
    "profesionales contra los Anexos 02 y 03. Eso es estandar de ORO "
    "(juicio humano); MIMIC es plata (codigos administrativos). "
    "Preprocesamiento: anonimizacion de identificadores en el texto "
    "libre, descarte de textos demasiado breves y 2,463 duplicados "
    "eliminados — sin deduplicar, el mismo texto caia en entrenamiento y "
    "prueba y las metricas salian infladas. Quedan 6,336 casos. "
    "Advertencia de comparabilidad (decirla ANTES de que la pregunten): "
    "el reporte describe un evento que quien escribio YA identifico; en "
    "la nota clinica hay que ENCONTRARLO. Problemas distintos; el del "
    "reporte es mas facil por construccion. No entrar en longitudes de "
    "texto.",

    # 12 · Preprocesamiento (1): correcciones
    "Las dos correcciones con efecto medido. BUG DEL PUNTO: el mapeo "
    "guardaba A04.7 (como la norma) y MIMIC guarda A047. El cruce fallaba "
    "en silencio: 411 hospitalizaciones. Corregido y cruzando por "
    "prefijo (respeta la jerarquia CIE-10): 109,714 — factor 267x. 'El "
    "Tier A nunca fallo conceptualmente; fallo por un punto.' BUG DEL "
    "COMODIN: el patron podia atravesar la nota entera — 'hemocultivo "
    "positivo' en la pagina 1 se unia con 'cateter' en la pagina 8: dos "
    "hechos verdaderos por separado, una conclusion falsa. Acotado a 100 "
    "caracteres (mismo parrafo): cae el 63.7% de detecciones, todas "
    "espurias. Elegancia del cuadro: una correccion SUBE la cifra y la "
    "otra la BAJA — ambas hacia la verdad: el Tier A estaba ciego, el "
    "Tier B alucinaba. Verificacion del diagnostico: los patrones SIN "
    "comodin quedaron identicos (13,189 -> 13,189, cero variacion); solo "
    "cayeron los que tenian .* — el defecto era el ALCANCE, no la "
    "especificidad.",

    # 13 · Preprocesamiento (2): representacion y particiones
    "Como se convierte el texto en numeros: TF-IDF con dos vistas. "
    "Palabra (1-2 gramas): capta conceptos como 'wound infection'. "
    "Caracter (3-5 letras): 'infec' captura infection/infected/infeccion "
    "— sobrevive a abreviaturas y errores del texto clinico. La logica "
    "TF-IDF: pesa mas lo frecuente en ESTA nota pero raro en el corpus — "
    "la intuicion del lector experto: 'reintubacion' informa, 'paciente' "
    "no. (Autores si preguntan: IDF de Karen Sparck Jones 1972; SVM de "
    "Vapnik/Cortes 1995; uso scikit-learn.) Y la decision de protocolo "
    "mas importante: PARTICION POR PACIENTE con GroupShuffleSplit y "
    "asercion de no solape — un mismo paciente aporta varias notas con "
    "vocabulario compartido; particionar por nota inflaria las metricas. "
    "El vectorizador se ajusta SOLO en entrenamiento.",

    # 14 · Modelamiento: dos etapas
    "Dos preguntas en serie. Etapa 1: hay evento adverso? — binaria, "
    "entrenada con 33,564 notas SIN evento (la clase negativa). Etapa 2: "
    "de que naturaleza? — 8 clases del Anexo 02. La clase negativa es lo "
    "que convierte un clasificador en un DETECTOR: le da el derecho a "
    "decir 'aqui no hay nada'. La consola es la prueba: seis textos "
    "triviales (un punto, 'ok', 'paciente estable'), seis abstenciones "
    "con margen negativo. Antes de esto, disparaba hasta ante un punto. "
    "Y la cascada es la honestidad: evaluar la Etapa 2 sobre positivos "
    "garantizados asume un detector perfecto; la cifra real es la de las "
    "dos etapas encadenadas — la muestro en resultados.",

    # 15 · Ejecucion: entorno y coste
    "La lamina que responde lo que pidio el docente. Comandos reales: "
    "entorno virtual AISLADO (para no comprometer el interprete del "
    "sistema) con PyTorch para CUDA. Hardware: laptop HP Victus, GTX 1650 "
    "de 4 GB — hardware de estudiante, deliberado: el DUA exige "
    "procesamiento local y demuestra transferibilidad a cualquier "
    "hospital. La tabla es la factura medida de los logs: ablacion del "
    "corpus completo casi 18 horas; cada fine-tuning de transformer 3-4 "
    "horas; los cuatro TF-IDF, 1 minuto 46 segundos. Total registrado: "
    "33 h 54 min. El contraste que importa (tarjeta granate): "
    "transformers 9 h 47 min para F1 0.354; TF-IDF 1 min 46 s para "
    "0.459. TRESCIENTAS TREINTA Y TRES veces mas computo, para perder.",

    # 16 · Ejecucion: consola real
    "Salidas LITERALES de los logs, no reconstruidas. La historia: la "
    "corrida grande murio al 78.4% tras 17 horas — ahi esta el mensaje "
    "del checkpoint — y se reanudo completando el resto en 32 minutos. "
    "Ejecutar un pipeline real es disenar para el fallo. Abajo, la "
    "verificacion del comodin: CON .* caen 84%; SIN comodin, cero "
    "variacion — la prueba de que el diagnostico era correcto. Cada "
    "experimento tiene comando, log y duracion en la bitacora publica. "
    "Ofrecer aqui: 'tengo el notebook abierto, puedo ejecutarlo en "
    "vivo'.",

    # 17 · Confusor de epoca (LAMINA ESTRELLA — 2 minutos)
    "El momento Zech anunciado. Version intermedia: AUC 0.973. "
    "Espectacular, publicable... y FALSO. Causa: el mapeo solo tenia "
    "CIE-10; EE.UU. adopto CIE-10 en 2015 y MIMIC va de 2008 a 2019 — "
    "toda hospitalizacion anterior quedaba negativa POR CONSTRUCCION. El "
    "modelo no aprendio eventos adversos: aprendio a distinguir la "
    "PLANTILLA documental de cada epoca. Prueba forense: su rasgo de "
    "mayor peso era 'rdwsd', un encabezado de laboratorio sin contenido "
    "clinico. Correccion: mapeo extendido a CIE-9 con equivalencias "
    "verificadas y emparejamiento por epoca A NIVEL DE NOTA — control "
    "exacto: 28.92% de era CIE-10 en positivos y en negativos, "
    "diferencia cero. Tras corregir, los rasgos importantes son "
    "clinicos: caida, revision, rechazo, reestenosis. AUC honesto: "
    "0.843. Leccion: un AUC de 0.973 puede sostenerse entero en una "
    "senal espuria — la pregunta no es cuanto acierta el modelo, sino "
    "POR QUE acierta.",

    # 18 · Resultados: deteccion
    "Detector final: sensibilidad 0.762, especificidad 0.770, AUC 0.843 "
    "— intervalos por bootstrap AGRUPADO POR PACIENTE (remuestrear notas "
    "sueltas viola la independencia y estrecha artificialmente los "
    "intervalos). El numero operativo: el VPP del test es 0.793, pero el "
    "test esta balanceado; reajustado por Bayes a la prevalencia real "
    "del 20.12%: VPP OPERATIVO 0.455. Significado: de cada 100 alertas, "
    "45 son eventos reales — frente a 20 que encuentra la revision al "
    "azar. Multiplica por 2.26 la eficiencia del revisor capturando el "
    "76% de los eventos. Rol declarado: FILTRO DE CRIBADO, no arbitro — "
    "no decide, prioriza. Reportar 0.455 y no 0.793 es honestidad: es la "
    "cifra con la que un servicio de calidad planifica su carga real.",

    # 19 · Resultados: ranking
    "Siete modelos, misma particion, misma semilla: comparables. Gana "
    "TF-IDF + LinearSVC con texto completo: F1-macro 0.459 en 48 "
    "segundos. El mejor transformer (ClinicalBERT ponderado): 0.354 en 4 "
    "horas de GPU. La fila MAS importante es la CUARTA: el mismo modelo "
    "campeon amputado al texto que ve BERT cae a 0.330 (-28%) con la "
    "misma arquitectura — eso AISLA la variable: el transformer no "
    "pierde por arquitectura, pierde porque su ventana de 256 tokens "
    "cubre el 9% del documento (mediana 3,148 tokens). Es pedir "
    "diagnostico a quien leyo la primera pagina de la historia. Y "
    "BioBERT (papers cientificos) queda ultimo: el lenguaje de los "
    "articulos no es el de las notas reales. Conclusion: ver todo el "
    "texto vale mas que entender mejor un fragmento; la mejora es "
    "ventana larga (Clinical-Longformer), no mas preentrenamiento.",

    # 20 · Resultados: cascada y juicio experto
    "Dos evaluaciones que casi nadie hace. CASCADA: Etapa 2 aislada da "
    "F1-micro 0.75; encadenada tras el detector real, 0.49 — menos "
    "34.5%. Reportar solo la aislada sobreestimaria un tercio; reporto "
    "ambas. JUICIO EXPERTO: 163 casos anotados con protocolo ciego "
    "(interfaz oculta el veredicto del sistema, orden aleatorio, tiempo "
    "cronometrado), 78 por duplicado con un segundo evaluador "
    "independiente. Contra el consenso: sensibilidad 0.945. Analisis "
    "ampliado con todos los casos: 0.914 con intervalo a la mitad. La "
    "especificidad contra experto es 0.54, coherente con cribado. Y el "
    "dato que cierra el circulo: el 25% de los casos que la codificacion "
    "declaraba negativos son eventos REALES para el experto — el "
    "subregistro, medido con mis propios datos. Si preguntan por que la "
    "especificidad solo sobre el consenso: el piloto se muestreo por "
    "estratos definidos por la salida del sistema; leerla del total "
    "seria sesgo de verificacion.",

    # 21 · Resultados: espanol
    "Ya no es promesa: cuatro tareas medidas sobre 6,336 casos peruanos "
    "con etiqueta de oro, prediccion fuera de pliegue (cada caso "
    "predicho por un modelo que nunca lo vio). Evento adverso vs "
    "incidente — la distincion central de la norma: 86% de exactitud. "
    "Naturaleza del Anexo 02: 85%. Severidad del Anexo 03: 73%. Y por "
    "primera vez el CODIGO ESPECIFICO de evento (41 clases modelables): "
    "76%. Ademas: los 8,799 casos estan en un Excel con codigo humano, "
    "prediccion y acierto/error marcado — revisable caso por caso. "
    "Hallazgo institucional: 'Gestion de la organizacion' era inviable "
    "en MIMIC (24 casos, F1 cero) y aqui alcanza F1 0.85 con 1,400 "
    "ejemplos — el corpus peruano rescata clases que el ingles no "
    "cubre.",

    # 22 · Posicion frente a la literatura (NUEVA — lo que pidio el docente)
    "Esta lamina responde directamente al pedido del docente: comparar mis "
    "metricas con las de los autores. Empiezo con la advertencia (esta en "
    "la bajada): los estandares difieren — se declara y se compara igual. "
    "Murff 2011 en JAMA: el PLN alcanza sensibilidades de 59 a 91% por "
    "complicacion, contra 5 a 46% de los indicadores por codigos — mi "
    "trabajo esta en el mismo regimen: 0.762 contra codigos, 0.945 contra "
    "experto. Classen 2011: eventos en el 33% de las admisiones y la "
    "notificacion voluntaria pierde el 90% — yo lo medi con datos propios: "
    "25% de los negativos eran reales. Zech 2018: AUC 0.931 que cae a "
    "0.815 fuera del hospital, la red identificaba el hospital en 99.9% — "
    "mi paralelo exacto: 0.973 espurio a 0.843 controlado. Y Li 2022: la "
    "ventana larga supera a ClinicalBERT — coherente con mi -28% al "
    "truncar. Cierre: 'la literatura reporta estos tres fenomenos por "
    "separado; este trabajo los replica y los CONTROLA con datos propios'.",

    # 23 · GEMSES destino
    "Ultima caja del diagrama; breve porque excede el alcance de PLN del "
    "curso. Lo detectado alimenta la matriz de priorizacion "
    "institucional — la vigente en EsSalud — que combina frecuencia e "
    "impacto y asigna banda (Verde/Amarillo/Rojo) y RESPONSABLE: "
    "servicio, departamento o direccion segun la banda. Es la "
    "diferencia entre un modelo que clasifica y un sistema que "
    "gestiona. (30-40 segundos maximo.)",

    # 24 · El aporte (NUEVA — 'vender el trabajo', pedido del docente)
    "La lamina de VENTA que pidio el docente. 'Mi aporte no es un algoritmo "
    "nuevo — es triple.' METODOLOGICO: un protocolo de auditoria de validez "
    "para supervision debil en texto clinico — siete modos de fallo "
    "identificados, medidos y corregidos; el confusor de epoca como caso "
    "ejemplar. EMPIRICO: el aislamiento experimental de la ventana de "
    "contexto como el factor que decide esta tarea — mismo algoritmo, "
    "mismos datos, -28% al truncar; eso define la via de mejora correcta. "
    "APLICADO: el primer acople documentado de PLN con la taxonomia "
    "normativa peruana — Anexos 02 y 03 — validado con etiqueta de oro en "
    "espanol. Remate: 'la literatura reporta metricas; este trabajo reporta "
    "metricas AUDITADAS — y las lleva por primera vez a la norma y al "
    "idioma de nuestro sistema de salud'.",

    # 25 · Conclusiones
    "Abrir con la numero UNO, que responde al problema original — es lo "
    "que el docente pidio explicitamente: '¿Resuelve el problema? La capa "
    "tecnica si: el sistema encuentra lo que nadie notifico — 25% de los "
    "negativos por codigo eran reales — y multiplica por 2.26 la "
    "eficiencia del revisor. La capa institucional queda LISTA: el piloto "
    "tiene sus parametros calculados.' Luego destacar dos mas: el modelo "
    "lexico supero al transformer con la causa MEDIDA (ventana de "
    "contexto, no arquitectura; la mejora es secuencia larga), y la "
    "AUDITORIA de validez como aporte central — siete modos de fallo, con "
    "el confusor de epoca como caso ejemplar; sin ella, las metricas "
    "habrian medido los defectos del corpus, no los eventos. Las demas "
    "estan escritas: "
    "cascada, espanol, destino institucional.",

    # 24 · Limitaciones
    "Decirlas con calma: SUMAN credibilidad. Las metricas sobre MIMIC "
    "miden acuerdo con codigos (estandar de plata), no verdad clinica — "
    "por eso existe la validacion experta. Esa validacion esta "
    "restringida a eventos de infeccion. El conjunto de prueba se "
    "reutilizo entre iteraciones — sesgo optimista no cuantificado, "
    "declarado; el umbral no se ajusto como mitigacion. Y el autor "
    "participo como anotador — mitigado con interfaz ciega, orden "
    "aleatorio, cronometraje y segundo evaluador independiente. Ninguna "
    "limitacion es oculta; todas tienen mitigacion documentada.",

    # 25 · Recomendaciones
    "Tres lineas, cada una anclada en un hallazgo. (1) Clinical-"
    "Longformer — ventana de 4,096 tokens: mi experimento demostro que "
    "el cuello de botella es el contexto; es una prediccion verificable. "
    "(2) Cuando exista acceso institucional: entrenar sobre notas "
    "clinicas peruanas con un encoder en espanol — el ERSP ya demostro "
    "que la taxonomia es aprendible en nuestro idioma. (3) Piloto "
    "institucional del cribado, con el VPP 0.455 como parametro de "
    "planificacion de carga.",

    # 26 · Referencias
    "34 fuentes en el articulo. Destacar solo: MIMIC-IV/PhysioNet, "
    "Ratner (Snorkel), Geirhos (shortcut learning), Zech (el caso "
    "clinico), Li 2022 (Clinical-Longformer) y la Directiva de EsSalud "
    "con los Anexos 02 y 03. JAMAS leerlas. (10 segundos.)",

    # 27 · Cierre
    "La idea que resume el trabajo: empece con un AUC de 0.973 y termine "
    "con 0.843 — y ese descenso es la mayor contribucion, porque el "
    "primero media un atajo y el segundo mide el fenomeno. En deteccion "
    "de eventos adversos, donde cada falso negativo es un dano invisible "
    "para la gestion, el aporte no es una metrica alta: es una metrica "
    "en la que se puede CONFIAR. El codigo es publico, la bitacora es "
    "auditable y el sistema ya habla espanol. Muchas gracias.",
]

pr = Presentation(P)
assert len(pr.slides) == len(NOTAS) == 29, (
    f"desfase: {len(pr.slides)} laminas vs {len(NOTAS)} notas")

for slide, nota in zip(pr.slides, NOTAS):
    slide.notes_slide.notes_text_frame.text = nota

pr.save(P)
print(f"[OK] {len(NOTAS)} notas del orador incrustadas en:")
print("    ", P)
